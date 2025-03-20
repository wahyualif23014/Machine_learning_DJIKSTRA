import streamlit as st
import pandas as pd
import docx
import PyPDF2
import pptx
import pytesseract
from PIL import Image
from io import BytesIO
from transformers import pipeline
import joblib
import re

summarizer = pipeline("summarization")

model = joblib.load("document_classifier.pkl")  

def read_text(file):
    return file.read().decode("utf-8")

def read_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def read_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def read_excel(file):
    df = pd.read_excel(file) if file.name.endswith(".xlsx") else pd.read_csv(file)
    return df.to_string()

def read_pptx(file):
    ppt = pptx.Presentation(file)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_image(file):
    image = Image.open(file)
    return pytesseract.image_to_string(image)



def summarize_text(text):
    return summarizer(text, max_length=100, min_length=30, do_sample=False)[0]["summary_text"]

def classify_document(text):
    predicted_category = model.predict([text])[0]
    return predicted_category

st.title("📂 AI File Analyzer & Summarizer")
st.write("Upload file dan AI akan menentukan jenisnya serta merangkum isinya!")

uploaded_file = st.file_uploader("Pilih file untuk dianalisis:", 
    type=["txt", "docx", "pdf", "csv", "xlsx", "pptx", "jpg", "png", "mp3", "wav"])

if uploaded_file is not None:
    file_type = uploaded_file.name.split(".")[-1]
    
    if file_type in ["txt"]:
        text = read_text(uploaded_file)
    elif file_type in ["docx"]:
        text = read_docx(uploaded_file)
    elif file_type in ["pdf"]:
        text = read_pdf(uploaded_file)
    elif file_type in ["csv", "xlsx"]:
        text = read_excel(uploaded_file)
    elif file_type in ["pptx"]:
        text = read_pptx(uploaded_file)
    elif file_type in ["jpg", "png"]:
        text = read_image(uploaded_file)
    else:
        st.error("Format file tidak didukung!")
        text = ""

    if text:
        st.subheader("📜 Isi File:")
        st.write(text[:1000])  
        document_type = classify_document(text)
        st.subheader(f"📌 Jenis Dokumen: {document_type}")

        st.subheader("📌 Ringkasan AI:")
        st.write(summarize_text(text)) 
# streamlit run resume.py
